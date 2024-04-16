from pptx import Presentation
from docx import Document
import PyPDF2
from google.cloud import translate_v3beta1 as translate
import boto3
from botocore.exceptions import BotoCoreError, NoCredentialsError
from sklearn.metrics.pairwise import cosine_similarity
import cohere
import numpy as np
import pandas as pd
import os
import requests
import uuid
import streamlit as st
from io import BytesIO

import os
import autogen
from autogen import ConversableAgent, UserProxyAgent, AssistantAgent, GroupChat, GroupChatManager
from autogen.coding import LocalCommandLineCodeExecutor
from autogen.oai.openai_utils import OAI_PRICE1K
from autogen.oai.client import OpenAIWrapper, ModelClient
from types import SimpleNamespace
import tempfile
from typing import Annotated, Literal, Any, Dict, List, Union
import inspect
from langchain_community.tools.tavily_search import TavilySearchResults
from langchain.utilities.tavily_search import TavilySearchAPIWrapper
import json
from openai import OpenAI


# App title
st.title('Document Translation and Cosine Similarity Comparison')

# File upload
input_file = st.file_uploader("Choose Input File", type=["txt", "pptx", "docx", "pdf"])
ground_truth_file = st.file_uploader("Choose Ground Truth File", type=["txt", "pptx", "docx", "pdf"])

# Dictionary of common languages and their codes for each translation service
languages_dict = {
    "English": {"DeepL": "EN", "AWS": "en", "Google": "en", "Azure": "en"},
    "Japanese": {"DeepL": "JA", "AWS": "ja", "Google": "ja", "Azure": "ja"},
    "French": {"DeepL": "FR", "AWS": "fr", "Google": "fr", "Azure": "fr"},
    "Spanish": {"DeepL": "ES", "AWS": "es", "Google": "es", "Azure": "es"},
    "German": {"DeepL": "DE", "AWS": "de", "Google": "de", "Azure": "de"},
    "Chinese": {"DeepL": "ZH", "AWS": "zh", "Google": "zh-CN", "Azure": "zh-Hans"},
    "Russian": {"DeepL": "RU", "AWS": "ru", "Google": "ru", "Azure": "ru"},
    "Portuguese": {"DeepL": "PT", "AWS": "pt", "Google": "pt", "Azure": "pt"},
    "Italian": {"DeepL": "IT", "AWS": "it", "Google": "it", "Azure": "it"},
    "Dutch": {"DeepL": "NL", "AWS": "nl", "Google": "nl", "Azure": "nl"},
    "Polish": {"DeepL": "PL", "AWS": "pl", "Google": "pl", "Azure": "pl"},
    "Turkish": {"DeepL": "TR", "AWS": "tr", "Google": "tr", "Azure": "tr"},
    "Korean": {"DeepL": "KO", "AWS": "ko", "Google": "ko", "Azure": "ko"},
    "Arabic": {"DeepL": "AR", "AWS": "ar", "Google": "ar", "Azure": "ar"},
    "Hebrew": {"DeepL": "HE", "AWS": "he", "Google": "iw", "Azure": "he"},
    "Hindi": {"DeepL": "HI", "AWS": "hi", "Google": "hi", "Azure": "hi"},
    "Swedish": {"DeepL": "SV", "AWS": "sv", "Google": "sv", "Azure": "sv"},
    # Add more languages as necessary
}

# Create a select box with language options
def create_language_selectbox(languages_dict):
    sorted_languages = sorted(languages_dict.keys())
    selected_language = st.selectbox("Select Output Language", sorted_languages)
    return selected_language, [languages_dict[selected_language].get(service, None) for service in ["DeepL", "AWS", "Google", "Azure"]]
target_lang, target_codes = create_language_selectbox(languages_dict)

# Txt extract
def extract_text_from_txt(file):
    return file.read().decode('utf-8')

# Pptx extract
def extract_text_from_pptx(file):
    prs = Presentation(file)
    extracted_text = []

    # Handle all shapes
    def handle_shape(shape):
        if hasattr(shape, "text"):
            return shape.text
        if shape.shape_type == 14:  # Placeholder for tables
            text = ''
            for row in shape.table.rows:
                for cell in row.cells:
                    text += cell.text + " "
            return text
        return ""

    # Handle grouped shapes
    def handle_grouped_shapes(shapes):
        grouped_text = ''
        for shape in shapes:
            if shape.shape_type == 6:  # Group shape
                grouped_text += handle_grouped_shapes(shape.shapes)
            else:
                grouped_text += handle_shape(shape)
        return grouped_text

    # Iterate through each slide
    for slide_number, slide in enumerate(prs.slides, start=1):
        slide_text = f"Slide {slide_number}:\n"
        # Iterate through each shape in each slide
        for shape in slide.shapes:
            if shape.shape_type == 6:  # Group shape
                slide_text += handle_grouped_shapes(shape.shapes)
            else:
                slide_text += handle_shape(shape)
        extracted_text.append(slide_text)

    # Merge extracted text
    merged_text = '\n'.join(text for text in extracted_text)

    return merged_text

# Cell extract
def extract_text_from_cell(cell):
    return '\n'.join(paragraph.text for paragraph in cell.paragraphs)

# Docx extract
def extract_text_from_docx(file):
    try:
        # Load the .docx file
        doc = Document(file)
        
        # Extract text from the document body
        full_text = [para.text for para in doc.paragraphs]

        # Extract text from headers and footers
        for section in doc.sections:
            header = section.header
            footer = section.footer
            
            # Extracting text from the header
            if header:
                full_text.extend([para.text for para in header.paragraphs])
            
            # Extracting text from the footer
            if footer:
                full_text.extend([para.text for para in footer.paragraphs])
        
        # Extract text from tables
        for table in doc.tables:
            for row in table.rows:
                for cell in row.cells:
                    full_text.append(extract_text_from_cell(cell))

        # Join
        return '\n'.join(full_text)
    
    except Exception as e:
        print(f"An error occurred: {e}")
        return None

# PDF extract
def extract_text_from_pdf(file):
    try:
        # Attempt to open the PDF file
        with open(file, 'rb') as file:
            # Create a PDF reader object
            pdf_reader = PyPDF2.PdfReader(file)
            
            # Initialize a variable to store the extracted text
            extracted_text = ''
            
            # Iterate through each page in the PDF
            for page in pdf_reader.pages:
                # Extract text from the page and add it to the accumulated text
                page_text = page.extract_text()
                if page_text:
                    extracted_text += page_text + '\n'
                else:
                    print(f"No text could be extracted from page {pdf_reader.pages.index(page)+1}.")
            
            return extracted_text
    except FileNotFoundError:
        print("The file was not found. Please check the path.")
        return None
    except PyPDF2.errors.PdfReadError:
        print("Error reading the PDF. It may be encrypted or corrupted.")
        return None
    except Exception as e:
        print(f"An error occurred: {e}")
        return None

# Select appropriate extract
def handle_file(file, file_type):
    if file_type == 'text/plain':
        text = extract_text_from_txt(file)
    elif file_type == 'application/vnd.openxmlformats-officedocument.presentationml.presentation':
        text = extract_text_from_pptx(file)
    elif file_type == 'application/vnd.openxmlformats-officedocument.wordprocessingml.document':
        text = extract_text_from_docx(file)
    elif file_type == 'application/pdf':
        text = extract_text_from_pdf(file)
    else:
        text = "Unsupported file type."
    return text

# Detect language
def detect_language(input_text):

    # Initialize client
    client = translate.TranslationServiceClient()

    # Project ID & location
    project_id = "translation-420503"
    location = 'global'  # The API currently only supports 'global' as the location
    parent = f"projects/{project_id}/locations/{location}"

    # Call API
    response = client.detect_language(
        request={
            "parent": parent,
            "content": input_text,
            "mime_type": "text/plain"
        }
    )

    # Extract detected language code
    if response.languages:
        source_lang = response.languages[0].language_code
        return source_lang
    else:
        print("No language detected.")
        return None
    
# Execute extract
def execute_extract(input_file, ground_truth_file):
    # Initialize default values for return variables
    input_text = gt_text = source_lang = None
    
    if input_file is not None:
        # Use the file details to determine the type
        file_details = input_file.type
        file_bytes = BytesIO(input_file.getvalue())
        
        # Extract text based on the file type
        input_text = handle_file(file_bytes, file_details)
        
        # Determine the source language
        source_lang = detect_language(input_text)

    if ground_truth_file is not None:
        # Use the file details to determine the type
        file_details = ground_truth_file.type
        file_bytes = BytesIO(ground_truth_file.getvalue())
        
        # Extract text based on the file type
        gt_text = handle_file(file_bytes, file_details)
        
    return input_text, gt_text, source_lang

# Deepl translate function
def translate_text_with_deepl(deepl_api_key, input_text, target_lang):
    # Set url
    url = "https://api-free.deepl.com/v2/translate" # Alternatively, "https://api.deepl.com/v2/translate"
    
    # Set data
    data = {
        'auth_key': deepl_api_key,
        'text': input_text,
        'target_lang': target_lang
    }
    
    # Call API
    try:
        response = requests.post(url, data=data)
        response.raise_for_status()
        
        # Extract translated text
        response_data = response.json()
        translated_text = response_data['translations'][0]['text']
        return translated_text
    except requests.exceptions.HTTPError as errh:
        print(f"HTTP Error: {errh}")
    except requests.exceptions.ConnectionError as errc:
        print(f"Error Connecting: {errc}")
    except requests.exceptions.Timeout as errt:
        print(f"Timeout Error: {errt}")
    except requests.exceptions.RequestException as err:
        print(f"Error: {err}")
    return None

# AWS translate function
def translate_text_with_amazon(input_text, source_lang, target_lang):
    try:
        # Initialize client
        client = boto3.client('translate', region_name='us-east-1')

        # Call API
        response = client.translate_text(
            Text=input_text,
            SourceLanguageCode=source_lang,
            TargetLanguageCode=target_lang
        )
        
        # Extract translated text
        return response['TranslatedText']
    except NoCredentialsError:
        return "Error: No credentials provided or found."
    except BotoCoreError as e:
        return f"An AWS error occurred: {e}"
    except Exception as e:
        return f"An error occurred: {e}"

# Google translate function
def translate_text_with_google(input_text, source_lang, target_lang):
    # Initialize client
    client = translate.TranslationServiceClient()

    # Project ID & location
    project_id = "translation-420503"
    location = 'global'  # The API currently only supports 'global' as the location
    parent = f"projects/{project_id}/locations/{location}"
    
    # Call API
    response = client.translate_text(
            request={
                "parent": parent,
                "contents": [input_text],
                "mime_type": 'text/plain',
                "source_language_code": source_lang,  # Can be set to 'auto'
                "target_language_code": target_lang
            }
        )
    
    # Extract translated text
    translated_text = ' '.join(translation.translated_text for translation in response.translations)

    return translated_text.strip()

# Azure translate function
def translate_text_with_azure(azure_subscription_key, input_text, source_lang, target_lang):
    # Construct url
    azure_endpoint = 'https://api.cognitive.microsofttranslator.com'
    path = '/translate?api-version=3.0'
    constructed_url = azure_endpoint + path

    # Set params
    params = {
        'from': source_lang,
        'to': target_lang
    }

    # Set headers
    headers = {
        'Ocp-Apim-Subscription-Key': azure_subscription_key,
        'Ocp-Apim-Subscription-Region': 'centralus',
        'Content-type': 'application/json',
        'X-ClientTraceId': str(uuid.uuid4())
    }

    # Set body
    body = [{'text': input_text}]

    # Call API
    response = requests.post(constructed_url, params=params, headers=headers, json=body)
    
    # Extract translated text
    if response.status_code == 200:
        response_data = response.json()
        translated_text = response_data[0]['translations'][0]['text']
        return translated_text
    else:
        print(f"Failed to translate text: HTTP {response.status_code}")
        return None

def cohere_emb_sim(docs, gt_text, names, display_table=True):
    # Initialize Cohere client
    co = cohere.Client(cohere_api_key)

    # Convert input and translated text to embeddings
    doc_emb = co.embed(texts=docs, model="embed-multilingual-v3.0", input_type="search_document", embedding_types=['float'])

    # Convert ground truth to an embedding
    gt_emb = co.embed(texts=[gt_text], model="embed-multilingual-v3.0", input_type="search_query", embedding_types=['float'])

    # Extract the embeddings
    gt_emb_array = np.array(gt_emb.embeddings.float_)
    doc_emb_array = np.array(doc_emb.embeddings.float_)

    # Compute the cosine similarity between the ground truth and document embeddings
    similarity_scores = cosine_similarity(gt_emb_array, doc_emb_array)

    # Since gt_embeddings_array contains only one set of embeddings, simplify the output
    similarity_scores = similarity_scores.flatten()  # This will make it a 1D array of length 5

    # Convert scores to a numpy array for indexing
    similarity_scores = np.array(similarity_scores)

    # Create a DataFrame
    df = pd.DataFrame({
        "Model": names,
        "Cosine Similarity": similarity_scores
    })
    
    # # Highlight the maximum similarity score
    # def highlight_max(s):
    #     is_max = s == s.max()
    #     return ['background-color: green' if v else '' for v in is_max]
    
    # # Display the DataFrame with highlighting
    # st.table(df.style.apply(highlight_max, subset=['Cosine Similarity']))
    
    if display_table:
        # Highlight the maximum similarity score
        def highlight_max(s):
            is_max = s == s.max()
            return ['background-color: green' if v else '' for v in is_max]
        
        # Display the DataFrame with highlighting
        st.table(df.style.apply(highlight_max, subset=['Cosine Similarity']))

# Input
with open('/Users/stephenchettiath/Library/Mobile Documents/com~apple~CloudDocs/Deep Learning/API Keys/azure_subscription_key.txt', 'r') as file:
    # Read the content of the file
    azure_subscription_key = file.read()
with open('/Users/stephenchettiath/Library/Mobile Documents/com~apple~CloudDocs/Deep Learning/API Keys/aws_access_key.txt', 'r') as file:
    aws_access_key = file.read()
with open('/Users/stephenchettiath/Library/Mobile Documents/com~apple~CloudDocs/Deep Learning/API Keys/aws_secret_key.txt', 'r') as file:
    aws_secret_key = file.read()
google_api_key = "/Users/stephenchettiath/Library/Mobile Documents/com~apple~CloudDocs/Deep Learning/API Keys/google_cloud_translation.json"
with open('/Users/stephenchettiath/Library/Mobile Documents/com~apple~CloudDocs/Deep Learning/API Keys/deepl_api_key.txt', 'r') as file:
    deepl_api_key = file.read()
with open('/Users/stephenchettiath/Library/Mobile Documents/com~apple~CloudDocs/Deep Learning/API Keys/cohere_api_key.txt', 'r') as file:
    cohere_api_key = file.read()
with open('/Users/stephenchettiath/Library/Mobile Documents/com~apple~CloudDocs/Deep Learning/API Keys/openai_api_key.txt', 'r') as file:
    openai_api_key = file.read().strip()
with open('/Users/stephenchettiath/Library/Mobile Documents/com~apple~CloudDocs/Deep Learning/API Keys/tavily_api_key.txt', 'r') as file:
    tavily_api_key = file.read().strip()

# Set environment variables
os.environ['AWS_ACCESS_KEY_ID'] = aws_access_key
os.environ['AWS_SECRET_ACCESS_KEY'] = aws_secret_key
os.environ["GOOGLE_APPLICATION_CREDENTIALS"] = google_api_key
os.environ['OPENAI_API_KEY'] = openai_api_key
os.environ['TAVILY_API_KEY'] = tavily_api_key

# Extract
input_text, gt_text, source_lang = execute_extract(input_file, ground_truth_file)

# Initialize variables to None
deepl_translated_text = aws_translated_text = google_translated_text = azure_translated_text = None

# Translate
if input_text is not None and gt_text is not None:
    deepl_translated_text = translate_text_with_deepl(deepl_api_key, input_text, target_codes[0])
    aws_translated_text = translate_text_with_amazon(input_text, source_lang, target_codes[1])
    google_translated_text = translate_text_with_google(input_text, source_lang, target_codes[2])
    azure_translated_text = translate_text_with_azure(azure_subscription_key, input_text, source_lang, target_codes[3])

# Docs
docs = [input_text,
            deepl_translated_text,
            aws_translated_text,
            google_translated_text,
            azure_translated_text]

# Names
names = ["Input Text", "DeepL", "AWS", "Google", "Azure"]

# Translate text
if st.button('Translate'):
    cohere_emb_sim(docs, gt_text, names)

# Create the prompt
prompt = f"""
Act as a linguistic synthesis expert. You have been given an input text and four machine translations of this text. Your task is to combine these translations into a single, coherent version. This version should integrate the nuances and strengths of each individual translation and sound natural in its native language. Address the common issues found in machine translations such as lack of context, poor handling of cultural nuances, struggles with idioms, technical jargon, and ambiguous language. Ensure the final output avoids any inappropriate translations, inaccuracies, grammatical errors, and confusing language. After creating the final output, refer back to the original input text to identify and implement any enhancements that could improve the quality of the final translation.

Here are the details:

Input Text: {input_text}
Translation 1: {deepl_translated_text}
Translation 2: {aws_translated_text}
Translation 3: {google_translated_text}
Translation 4: {azure_translated_text}

Please synthesize these translations into a single, enhanced version that captures the best elements of each, and provide reasoning for key integration choices made during the process.
"""
# Define a session state variable to hold the result if it doesn't already exist
if 'translation_result' not in st.session_state:
    st.session_state.translation_result = ""

# Display the text area with the value from the session state
st.text_area("GPT-4 Turbo Result", value=st.session_state.translation_result, height=150)

# Enhance with GPT-4 turbo
#st.text_area("GPT-4 Turbo Result", height=150)
if st.button('Enhance w/GPT-4 Turbo'):
    client = OpenAI()
    response = client.chat.completions.create(
    model="gpt-4-turbo-2024-04-09",
    messages=[
        {"role": "system", "content": "You are a helpful assistant."},
        {"role": "user", "content": prompt}
    ]
    )
    # result = response.choices[0].message.content
    # st.write(result)
    # Extract the response content and update the session state
    st.session_state.translation_result = response.choices[0].message.content
    # Redisplay the page with the new result in the text area
    st.experimental_rerun()
    
# Docs
docs = [input_text, deepl_translated_text, aws_translated_text, google_translated_text, azure_translated_text, st.session_state.translation_result]

# Names
names = ["Input Text", "DeepL", "AWS", "Google", "Azure", "GPT-4 Turbo Enhanced"]

# Rerun similarity
if st.button('Recalculate Cosine Similarities'):
    cohere_emb_sim(docs, gt_text, names)

# Enhance with group chat
def run_group_chat(prompt):
    # Search engine tool
    search = TavilySearchAPIWrapper()
    tavily_tool = TavilySearchResults(api_wrapper=search)
    search_tool = TavilySearchResults()

    def generate_llm_config(tool):
        # Define the function schema based on the tool's args_schema
        function_schema = {
            "name": tool.name.lower().replace(" ", "_"),
            "description": tool.description,
            "parameters": {
                "type": "object",
                "properties": {},
                "required": [],
            },
        }

        if tool.args is not None:
            function_schema["parameters"]["properties"] = tool.args

        return function_schema

    # Create the calculator function
    Operator = Literal["+", "-", "*", "/"]
    def calculator(a: int, b: int, operator: Annotated[Operator, "operator"]) -> int:
        if operator == "+":
            return a + b
        elif operator == "-":
            return a - b
        elif operator == "*":
            return a * b
        elif operator == "/":
            return int(a / b)
        else:
            raise ValueError("Invalid operator")
    
    
    # Create the gpt conversable agent
    gpt = ConversableAgent(
        "gpt",
        system_message="Act as a linguistic synthesis expert who thoroughly looks for ways to enhance translations, ensuring that every translation integrates the subtleties, context, and cultural nuances of the source text. You strive to make every individual translation sound natural in its native language while addressing the common issues found in machine translations such as lack of context, poor handling of cultural nuances, struggles with idioms, technical jargon, and ambiguous language. Confirm that the final output avoids any inappropriate translations, inaccuracies, grammatical errors, and confusing language. You have been given a search engine tool to perform web-scale search and a simple calculator to perform simple calculations to solve tasks. Reply TERMINATE when the task is done.",
        llm_config={"functions": [generate_llm_config(search_tool)], "config_list": [{"model": "gpt-4-turbo-2024-04-09", "temperature": 0.7, "api_key": os.environ.get("OPENAI_API_KEY"), "timeout": 120}]},
        code_execution_config=False,
        function_map=None,
        human_input_mode="NEVER", # Alternatively, "ALWAYS" or "TERMINATE"
    )

    # Register the calculator tool function with the agent
    gpt.register_for_execution(name="calculator")(calculator)

    # Register the search engine tool function with the agent
    gpt.register_function(
        function_map={
            search_tool.name: search_tool._run,
        }
    )

    # Create the other_gpt conversable agent
    other_gpt = ConversableAgent(
        "other_gpt",
        system_message="Act as a linguistic synthesis expert who thoroughly looks for ways to enhance translations, ensuring that every translation integrates the subtleties, context, and cultural nuances of the source text. You strive to make every individual translation sound natural in its native language while addressing the common issues found in machine translations such as lack of context, poor handling of cultural nuances, struggles with idioms, technical jargon, and ambiguous language. Confirm that the final output avoids any inappropriate translations, inaccuracies, grammatical errors, and confusing language. You have been given a search engine tool to perform web-scale search and a simple calculator to perform simple calculations to solve tasks. Reply TERMINATE when the task is done.",
        llm_config={"functions": [generate_llm_config(search_tool)], "config_list": [{"model": "gpt-4-turbo-2024-04-09", "temperature": 0.7, "api_key": os.environ.get("OPENAI_API_KEY"), "timeout": 120}]},
        code_execution_config=False,
        function_map=None,
        human_input_mode="NEVER", # Alternatively, "ALWAYS" or "TERMINATE"
    )

    # Register the calculator tool function with the agent
    other_gpt.register_for_execution(name="calculator")(calculator)

    # Register the search engine tool function with the agent
    other_gpt.register_function(
        function_map={
            search_tool.name: search_tool._run,
        }
    )

    # Initialize a group chat object
    group_chat = GroupChat(
        agents=[gpt, other_gpt],
        messages=[],
        max_round=10,
        send_introductions=True,
    )
    
    # Describe the agents for the group manager
    gpt.description = "An expert in linguistic synthesis who can also utilize a search engine tool and simple calculator to solve tasks. Powered by GPT-4 Turbo."
    other_gpt.description = "An expert in linguistic synthesis who can also utilize a search engine tool and simple calculator to solve tasks. Powered by GPT-4 Turbo."

    # Create the non-speaking group chat manager
    group_chat_manager = GroupChatManager(
        groupchat=group_chat,
        llm_config={"config_list": [{"model": "gpt-3.5-turbo-0125", "api_key": os.environ["OPENAI_API_KEY"]}]},
    )

    # Initiate the chat and return a summary of the whole conversation at the end
    chat_result = group_chat_manager.initiate_chat(
        gpt,
        message=prompt,
        summary_method="reflection_with_llm",
        max_turns=5,
    )
    
    return chat_result.summary

# Define a session state variable to hold the result if it doesn't already exist
if 'group_chat_result' not in st.session_state:
    st.session_state.group_chat_result = ""

# Display the text area with the value from the session state
st.text_area("Group Chat Summary", value=st.session_state.group_chat_result, height=150)

#st.text_area("Group Chat Summary", height=150)
if st.button('Run Group Chat'):
    # summary = run_group_chat(prompt)
    # st.write("Summary of the Group Chat:")
    # st.write(summary)
    # Extract the response content and update the session state
    st.session_state.group_chat_result = run_group_chat(prompt)
    # Redisplay the page with the new result in the text area
    st.experimental_rerun()